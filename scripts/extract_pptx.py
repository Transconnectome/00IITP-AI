import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_pptx(pptx_path, output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    prs = Presentation(pptx_path)
    content = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        image_count = 0
        
        # Extract notes if any
        notes_text = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                image_filename = f"slide_{i+1}_image_{image_count}.png"
                image_path = os.path.join(output_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(shape.image.blob)
                slide_text.append(f"[IMAGE: {image_filename}]")
        
        slide_content = f"--- Slide {i+1} ---\n"
        slide_content += "\n".join(slide_text)
        if notes_text:
            slide_content += f"\n\nNotes: {notes_text}"
        content.append(slide_content)
        
    return "\n\n".join(content)

if __name__ == "__main__":
    pptx_file = "docs/IdeaDeck.pptx"
    assets_dir = "docs/ideadeck_assets"
    extracted_text = extract_pptx(pptx_file, assets_dir)
    
    with open("docs/ideadeck_extracted.txt", "w") as f:
        f.write(extracted_text)
    
    print(f"Extracted content to docs/ideadeck_extracted.txt and images to {assets_dir}")
